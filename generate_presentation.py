import os
import sys
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6] # Blank layout

    # Colors
    bg_color = RGBColor(253, 251, 247)       # Warm Cream / Ivory
    navy_color = RGBColor(26, 46, 64)        # Deep Navy (Titles)
    taupe_color = RGBColor(156, 140, 112)    # Warm Taupe (Accent / Lines)
    charcoal_color = RGBColor(44, 44, 44)    # Charcoal (Body Text)
    sage_color = RGBColor(110, 138, 117)     # Sage Green (Subtitles)

    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_fit_picture(slide, img_path, left, top, max_w, max_h):
        if not os.path.exists(img_path):
            print(f"Warning: Image {img_path} not found.")
            return
        
        try:
            im = Image.open(img_path)
            im_w, im_h = im.size
            aspect = im_w / im_h
            
            w = max_w
            h = w / aspect
            if h > max_h:
                h = max_h
                w = h * aspect
                
            x = left + (max_w - w) / 2
            y = top + (max_h - h) / 2
            
            slide.shapes.add_picture(img_path, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
            
            # Draw a subtle borders/frame around the image
            border = slide.shapes.add_shape(
                1, # MSO_SHAPE.RECTANGLE
                Inches(x), Inches(y), Inches(w), Inches(h)
            )
            border.fill.background()
            border.line.color.rgb = taupe_color
            border.line.width = Pt(1)
        except Exception as e:
            print(f"Error adding picture {img_path}: {e}")

    def add_title_slide(title_text, subtitle_text, author_text):
        slide = prs.slides.add_slide(blank_slide_layout)
        apply_background(slide)

        # Title + Subtitle container
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        # Title
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Georgia'
        p.font.size = Pt(64)
        p.font.bold = True
        p.font.color.rgb = navy_color
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(24)

        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = 'Georgia'
        p2.font.size = Pt(20)
        p2.font.color.rgb = sage_color
        p2.alignment = PP_ALIGN.CENTER
        p2.space_after = Pt(40)

        # Author / Footer
        p3 = tf.add_paragraph()
        p3.text = author_text
        p3.font.name = 'Garamond'
        p3.font.size = Pt(14)
        p3.font.color.rgb = taupe_color
        p3.alignment = PP_ALIGN.CENTER

        # Accent Line
        shape = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE
            Inches(5.666), Inches(4.3), Inches(2.0), Inches(0.02)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = taupe_color
        shape.line.fill.background()

    def add_content_slide(title_text, bullet_points, img_path=None):
        slide = prs.slides.add_slide(blank_slide_layout)
        apply_background(slide)

        # Title Box
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        p = tf_title.paragraphs[0]
        p.text = title_text
        p.font.name = 'Georgia'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = navy_color

        # Elegant divider line
        line = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE
            Inches(1.0), Inches(1.6), Inches(11.333), Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = taupe_color
        line.line.fill.background()

        # Decide widths based on image availability
        text_width = 5.6 if img_path else 11.333

        # Content Box
        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(text_width), Inches(4.8))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_left = tf_content.margin_top = tf_content.margin_right = tf_content.margin_bottom = 0

        first = True
        for point in bullet_points:
            if first:
                p_bullet = tf_content.paragraphs[0]
                first = False
            else:
                p_bullet = tf_content.add_paragraph()
            
            p_bullet.space_after = Pt(14)
            p_bullet.level = 0
            
            # Format bold start if any
            if "**" in point:
                parts = point.split("**")
                for i, part in enumerate(parts):
                    run = p_bullet.add_run()
                    run.text = part
                    run.font.name = 'Garamond'
                    run.font.size = Pt(16 if img_path else 18)
                    run.font.color.rgb = charcoal_color
                    if i % 2 == 1:
                        run.font.bold = True
                        run.font.color.rgb = navy_color
            else:
                run = p_bullet.add_run()
                run.text = point
                run.font.name = 'Garamond'
                run.font.size = Pt(16 if img_path else 18)
                run.font.color.rgb = charcoal_color

        # Add image if path is provided
        if img_path:
            add_fit_picture(slide, img_path, left=7.0, top=2.0, max_w=5.3, max_h=4.5)

        # Elegant subtle footer
        footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.9), Inches(11.333), Inches(0.4))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "NetVigil — Observability & Threat Intelligence Platform"
        p_foot.font.name = 'Garamond'
        p_foot.font.size = Pt(10)
        p_foot.font.italic = True
        p_foot.font.color.rgb = taupe_color
        p_foot.alignment = PP_ALIGN.RIGHT

    # Slide 1: Title
    add_title_slide(
        "NetVigil",
        "Платформа активного моніторингу безпеки та інтелектуального аналізу загроз",
        "Роман Курмаш | Дипломний проект"
    )

    # Slide 2: Goal
    add_content_slide(
        "Мета проекту",
        [
            "**Автоматизація SecOps**: Перехід від пасивного спостереження до активного аналізу та автоматичного реагування на інциденти.",
            "**ШІ-аналіз загроз**: Залучення локальних великих мовних моделей (LLM) через Ollama для інтелектуального аналізу поведінки контейнерів.",
            "**Життєвий цикл інцидентів**: Об'єднання розрізнених повідомлень у єдиний державний реєстр інцидентів з відстеженням станів (Firing/Resolved).",
            "**Контекстуалізація подій**: Забезпечення повного аналітичного контексту для SecOps-аналітиків за рахунок крос-сервісної кореляції логів."
        ]
    )

    # Slide 3: Architecture
    add_content_slide(
        "Загальна архітектура системи",
        [
            "**Рівень обсервабіліті (Metrics & Logs)**: Promtail збирає логи контейнерів Docker та пересилає їх до Loki. Prometheus збирає метрики.",
            "**Рівень аналітики (AI-Adapter)**: Потоковий сканер логів вичитує нові записи з Loki, класифікує їх за допомогою LLM (Llama-3) та розраховує бал ризику.",
            "**Рівень сповіщень (Alertmanager & Bot)**: Alerts з високим ризиком надсилаються до Alertmanager, який тригерить Bot-мікросервіс для надсилання сповіщень.",
            "**SecOps Dashboard (Web-App)**: Веб-інтерфейс для візуалізації реєстру інцидентів, запуску ручної кореляції та ШІ-ретроспектив."
        ],
        img_path="documentation/diploma/images/architecture_diagram.png"
    )

    # Slide 4: Active Observability
    add_content_slide(
        "Концепція активної обсервабіліті",
        [
            "**Автоматичне виявлення сервісів**: Динамічне налаштування Promtail для збору логів з усіх запущених Docker-контейнерів.",
            "**Попередження помилок**: Спеціальні правила Prometheus/Alertmanager для виявлення критичних помилок, перевантаження пам'яті та мережевих аномалій.",
            "**Loki як єдине джерело**: Централізоване зберігання логів з можливостю швидкого пошуку за мітками контейнерів."
        ]
    )

    # Slide 5: AI Threat Detection
    add_content_slide(
        "ШІ-аналіз логів та виявлення загроз",
        [
            "**Потокове сканування**: Робота мікросервісу AI-Adapter у реальному часі для виявлення аномальних активностей.",
            "**Локальний LLM Llama-3**: Використання Ollama для збереження конфіденційності даних (без надсилання логів у хмарні API).",
            "**Аналіз ризиків**: Розрахунок Risk Score (0-10) та визначення типу загрози (SQL Injection, RCE, brute-force тощо).",
            "**Розумна фільтрація**: Надсилання до Alertmanager лише критичних подій (Risk >= 7) для запобігання втомі від оповіщень."
        ],
        img_path="documentation/diploma/images/log_analysis_flowchart.png"
    )

    # Slide 6: Stateful Incident Registry
    add_content_slide(
        "Реєстр подій та життєвий цикл інцидентів",
        [
            "**Агрегація подій**: Автоматичне групування сповіщень за парою (alert_name, container) у єдину сутність інциденту.",
            "**Життєвий цикл (Firing -> Resolved)**: Автоматичне оновлення стану інциденту при отриманні сигналу про усунення проблеми.",
            "**Хронологічний таймлайн**: Збереження всіх проміжних повідомлень та логів, що викликали інцидент, в межах одного запису.",
            "**Збереження даних**: Локальна база даних інцидентів у форматі JSON для швидкого доступу та аналізу."
        ],
        img_path="documentation/diploma/images/screenshot_telegram_alert.png"
    )

    # Slide 7: Cross-Container Correlation
    add_content_slide(
        "Крос-контейнерна кореляція логів",
        [
            "**Визначення взаємозв'язків**: Пошук спільних логів помилок на інших контейнерах у часовому вікні ±5 хвилин від початку інциденту.",
            "**Швидка діагностика**: Можливість миттєво побачити, чи була мережева атака на веб-сервер пов'язана із помилками в базі даних.",
            "**Спільні алерти**: Автоматичне виявлення інших активних сповіщень у той самий проміжок часу для побудови повної картини атаки."
        ]
    )

    # Slide 8: Automated AI Retrospectives
    add_content_slide(
        "Автоматичні ШІ-ретроспективи (Post-Mortem)",
        [
            "**Генерація post-mortem**: Автоматичний аналіз інциденту за допомогою локальної Llama-3 після його завершення.",
            "**Формат звіту (JSON)**: Звіт містить структурований аналіз першопричини (Root Cause), хронології (Timeline), наслідків (Impact) та рекомендацій.",
            "**Запобігання рецидивам**: Формування списку рекомендацій щодо покращення безпеки інфраструктури (Preventative Actions).",
            "**Україномовний аналіз**: Звіт формується державною мовою для спрощення SecOps-документації."
        ]
    )

    # Slide 9: UI Dashboard
    add_content_slide(
        "SecOps Dashboard та управління",
        [
            "**Інтерактивна таблиця**: Візуальний статус активності, рівні ризику, підсвічування статусів.",
            "**Двоколонковий Grid-інтерфейс деталей**: Ліва колонка містить історію повідомлень та сирі логи з Loki. Права колонка містить картки ручного запуску ШІ-аналізу та лог-кореляції.",
            "**Асинхронні запити**: Плавне завантаження результатів аналітики без перезавантаження сторінки."
        ],
        img_path="documentation/diploma/images/screenshot_dashboard.png"
    )

    # Slide 10: Conclusion
    add_content_slide(
        "Висновки та результати роботи",
        [
            "**Підвищення прозорості**: SecOps отримав єдиний пульт керування інцидентами з повною видимістю логів та хронології.",
            "**ШІ як асистент**: Локальний ШІ успішно знімає рутинне навантаження з аналітика, пропонуючи готові кроки вирішення.",
            "**Безпека даних**: Завдяки Ollama конфіденційні логи системи не залишають внутрішнього периметру інфраструктури.",
            "**Ефективність**: Шійдість розслідування інцидентів зросла завдяки автоматичній крос-контейнерній кореляції логів."
        ]
    )

    prs.save('NetVigil_Presentation.pptx')
    print("Presentation saved successfully as NetVigil_Presentation.pptx with images.")

if __name__ == '__main__':
    create_presentation()
